import os
import sys

try:
    from pptx import Presentation
except ImportError:
    print("python-pptx is not installed")
    sys.exit(1)

file_path = "cdCTF-Pitch-Deck.pptx"
if not os.path.exists(file_path):
    print("Pitch deck not found!")
    sys.exit(1)

prs = Presentation(file_path)

replacements = {
    "O'zbekiston uchun bepul,": "O'zbekiston uchun innovatsion,",
    "Bepul, amaliy, uch tilli": "Innovatsion, amaliy, uch tilli",
    "Bepul — bu strategiya,": "Keng qamrovli strategiya,",
    "Bepul ta'lim iste'dodlar": "Sifatli ta'lim iste'dodlar",
    "Bepul ta'lim — ko'p o'quvchi": "Ochiq ta'lim — ko'p o'quvchi",
    "O'quvchi uchun har doim bepul": "10,000$ sarmoya — 10 oylik o'sish uchun",
    "xayriya emas": "kengayuvchi B2B biznes modelidir",
    "bepul kiberxavfsizlik ta'limi": "innovatsion kiberxavfsizlik ta'limi",
    "o'quvchi bepul.": "juda qulay shartlarda jalb etiladi.",
    "bepul.": "ochiq.",
    "bepul": "innovatsion",
    "Bepul": "Innovatsion",
}

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        for paragraph in shape.text_frame.paragraphs:
            # Reconstruct the full text of the paragraph
            full_text = "".join(run.text for run in paragraph.runs)
            
            replaced = False
            for old_text, new_text in replacements.items():
                if old_text in full_text:
                    full_text = full_text.replace(old_text, new_text)
                    replaced = True
            
            if replaced:
                # Put everything in the first run to update the text, clear the rest.
                # This may lose intra-paragraph formatting (like one bold word), but guarantees the text is replaced.
                if len(paragraph.runs) > 0:
                    paragraph.runs[0].text = full_text
                    for i in range(1, len(paragraph.runs)):
                        paragraph.runs[i].text = ""

# Additionally, let's look for the last slide (Slide 11) and try to add a small note about the $10,000 investment if it's not already there.
# The slide with "Bir jamoa. Bir missiya."
for shape in prs.slides[-1].shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            full_text = "".join(run.text for run in paragraph.runs)
            if "biznes model bilan" in full_text:
                if len(paragraph.runs) > 0:
                    paragraph.runs[0].text = full_text + "\nSarmoya maqsadi: $10,000 (10 oylik jamoa, marketing va server xarajatlari uchun)"

out_path = "cdCTF-Pitch-Deck-Investor.pptx"
prs.save(out_path)
print(f"Successfully created {out_path}")
