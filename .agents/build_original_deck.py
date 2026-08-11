from pptx import Presentation
from pptx.util import Inches, Pt
import os
import shutil

# Assumes screenshots are already copied as screenshot_hero.png etc. from previous step.
# If not, let's copy them just in case.
base_artifact_dir = "/home/bozkurt/.gemini/antigravity/brain/b80c17e8-2aa7-44de-8fcc-5bd05055c753/"
try:
    if not os.path.exists('screenshot_hero.png'):
        shutil.copyfile(os.path.join(base_artifact_dir, [f for f in os.listdir(base_artifact_dir) if f.startswith('screenshot_hero')][0]), 'screenshot_hero.png')
    if not os.path.exists('screenshot_modules.png'):
        shutil.copyfile(os.path.join(base_artifact_dir, [f for f in os.listdir(base_artifact_dir) if f.startswith('screenshot_modules')][0]), 'screenshot_modules.png')
    if not os.path.exists('screenshot_ctf.png'):
        shutil.copyfile(os.path.join(base_artifact_dir, [f for f in os.listdir(base_artifact_dir) if f.startswith('screenshot_ctf')][0]), 'screenshot_ctf.png')
except Exception as e:
    print("Screenshot copies failed, maybe already present", e)

prs = Presentation('cdCTF-Pitch-Deck.pptx')

# Text replacements to make it Investor-ready, preserving original 100/100 layout
replacements = {
    "O'quvchi uchun bepul": "Talent Pipeline va B2B Ta'lim",
    "bepul ta'limni": "kiberxavfsizlik ta'limini",
    "bepul ta'lim": "ochiq ta'lim",
    "Homiyli CTF": "Korporativ Ta'lim",
    "Homiyli mahalliy CTF": "B2B Korporativ Ta'lim",
    "Kompaniyalar brendlangan musobaqa": "Kompaniyalar xodimlarini o'qitish",
    "Kompaniyalar brendlangan musobaqa va tadbirlarga homiylik qiladi — analitika bilan.": "Kompaniyalar IT xodimlarini xavfsiz kod yozish va infratuzilma himoyasiga o'qitadi.",
    "Homiylar": "Investitsiya Maqsadi",
    "Brendlangan CTF va fond hamkorligi": "Sarmoya: $10,000 (Jamoa, Marketing, Server)",
    "Bepul ta'limni ta'minlash": "Sifatli ta'limni ta'minlash",
    "Fond & grantlar": "Talent Pipeline",
    "Bepul ta'limni ta'minlash uchun grant, homiylik va davlat/xalqaro fondlar.": "Kompaniyalarga tasdiqlangan mutaxassislarni tavsiya qilish orqali HR-komissiyasi.",
    "O'zbekiston uchun bepul,": "O'zbekiston uchun ilg'or,",
    "xayriya emas": "B2B korporativ biznes"
}

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                full_text = "".join(run.text for run in paragraph.runs)
                replaced = False
                for k, v in replacements.items():
                    if k in full_text:
                        full_text = full_text.replace(k, v)
                        replaced = True
                
                if replaced:
                    if len(paragraph.runs) > 0:
                        # Retain the run formatting of the first run as much as possible
                        paragraph.runs[0].text = full_text
                        for i in range(1, len(paragraph.runs)):
                            paragraph.runs[i].text = ""

# Add the 3 screenshots carefully as new slides at the end so it doesn't mess up their layout
# We will use the blank layout (usually layout 6)
layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

screenshots = [
    ("cdCTF - Asosiy Oyna (Platforma Jonli)", "screenshot_hero.png"),
    ("cdCTF - O'quv Dasturi", "screenshot_modules.png"),
    ("cdCTF - Amaliyot va CTF Poligoni", "screenshot_ctf.png")
]

for title, img in screenshots:
    if os.path.exists(img):
        slide = prs.slides.add_slide(layout)
        
        # Add title text box at the top
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(36)
        try:
            # Try to match a nice dark color for the title if possible
            from pptx.dml.color import RGBColor
            p.font.color.rgb = RGBColor(30, 40, 50)
        except:
            pass
        
        # Center the 16:9 screenshot beautifully with some padding
        # 16 inches wide, 9 inches high. Let's make image 14.2 wide x 8 high.
        # Centered: Left = (16-14.2)/2 = 0.9, Top = (9-8)/2 = 0.5 (but we have title, so Top=1.0, Height=7.5, Width=13.3)
        slide.shapes.add_picture(img, Inches(1.3), Inches(1.0), Inches(13.33), Inches(7.5))

prs.save('cdCTF-Pitch-Deck-Final-OriginalDesign.pptx')
print("Saved cdCTF-Pitch-Deck-Final-OriginalDesign.pptx")
