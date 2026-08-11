from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

prs = Presentation()
# Set 16:9 aspect ratio
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

blank_slide_layout = prs.slide_layouts[6] 

def add_slide_with_image(title, subtitle, image_path=None):
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 15, 20)
    
    if image_path and os.path.exists(image_path):
        # Place image on the right side
        slide.shapes.add_picture(image_path, Inches(8.5), Inches(0), Inches(7.5), Inches(9))
        text_width = Inches(7)
    else:
        text_width = Inches(14)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), text_width, Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(54)
    p.font.name = "Arial"
    p.font.color.rgb = RGBColor(0, 255, 150)
    
    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), text_width, Inches(3.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(28)
    p2.font.name = "Arial"
    p2.font.color.rgb = RGBColor(220, 220, 220)
    return slide

# 1. Cover
add_slide_with_image("cdCTF", "O'zbekistonning innovatsion kiberxavfsizlik akademiyasi va Talent Pipeline platformasi.\n\nInvestitsiya maqsadi: $10,000", "/home/bozkurt/.gemini/antigravity/brain/b80c17e8-2aa7-44de-8fcc-5bd05055c753/cyber_cover_1786357013515.png")

# 2. Problem
add_slide_with_image("MUAMMO", "Kiberxavfsizlik kadrlarining keskin yetishmovchiligi:\n\n• O'zbek tilida amaliy, chuqurlashtirilgan ta'lim yo'q.\n• Iqtidorli yoshlar bor, lekin ishga joylashish tizimi yo'q.\n• Kompaniyalar xavfsizlik mutaxassislarini topishga qiynalmoqda.", "/home/bozkurt/.gemini/antigravity/brain/b80c17e8-2aa7-44de-8fcc-5bd05055c753/cyber_problem_1786357023887.png")

# 3. Solution
add_slide_with_image("YECHIM: cdCTF", "O'rganishdan birinchi ishgacha bo'lgan yagona ekotizim:\n\n• Akademiya: 8 modul, 64 dars (uz/ru/en).\n• CTF Platformasi: 97+ amaliy laboratoriyalar.\n• Karyera: Bitiruvchilarni to'g'ridan-to'g'ri kompaniyalarga ulash.", "/home/bozkurt/.gemini/antigravity/brain/b80c17e8-2aa7-44de-8fcc-5bd05055c753/cyber_solution_1786357033626.png")

# 4. Tech
add_slide_with_image("TEXNOLOGIYA VA XAVFSIZLIK", "Sanoat standartlaridagi kuchli arxitektura:\n\n• Zamonaviy Stack: React, Vite, Express 5, PostgreSQL, Drizzle\n• Xavfsiz Laboratoriyalar: Maxsus Docker konteynerlar orqali izolyatsiya qilingan (Sandboxed) arxitektura.\n• 0 Zaiflik: Qat'iy xavfsizlik auditi va 22 ta avtomatik testlardan o'tgan tizim.")

# 5. Business Model
add_slide_with_image("BIZNES MODEL (B2B)", "Aniq monetizatsiya rejalari:\n\n1. Talent Pipeline (HR Xizmati):\nKompaniyalarga amaliyotda sinalgan mutaxassislarni yetkazib berish va komissiya olish.\n\n2. Korporativ Ta'lim:\nKompaniyalarning o'z dasturchilarini xavfsizlikka o'qitish uchun yopiq B2B modullar.\n\n3. Premium Modullar:\nChuqurlashtirilgan (Advanced) topshiriqlar va xalqaro darajadagi maxsus sertifikatsiya xizmatlari.")

# 6. Traction
add_slide_with_image("TRACTION VA YUTUQLAR", "Loyiha ishlayapti va foydalanuvchilar o'smoqda:\n\n• Holat: MVP to'liq ishga tushirilgan (cdctf.uz)\n• Kontent Bazasi: 8 ta to'liq modul va 97 ta amaliy masala yaratib bo'lingan.\n• Multilingual: O'zbek, Rus va Ingliz tillarida parallel ishlash.\n• Foydalanuvchilar: Ilk qiziqish bildirgan yoshlar platformadan faol foydalanishni boshlagan.")

# 7. Ask
add_slide_with_image("INVESTITSIYA MAQSADI", "So'ralayotgan sarmoya: $10,000 (10 oylik Run-way)\n\nSarflanish rejasi:\n• Jamoani Kengaytirish (40%): Yangi murakkab laboratoriyalar yozadigan muhandislar jalb qilish.\n• Marketing (40%): Keng ommani jalb qilish va kuchli hamjamiyat qurish uchun reklamalar.\n• Infratuzilma (20%): CTF musobaqalari uchun bulutli (Cloud) serverlarni ta'minlash.")

# 8. Contact
slide_last = prs.slides.add_slide(blank_slide_layout)
background = slide_last.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(15, 15, 20)
txBox_last = slide_last.shapes.add_textbox(Inches(3), Inches(3.5), Inches(10), Inches(4))
tf_last = txBox_last.text_frame
p_last = tf_last.paragraphs[0]
p_last.text = "Bir jamoa. Bir missiya.\n\nHikmatbek Xudoyberganov\ncdctf.uz | @cdctf"
p_last.font.bold = True
p_last.font.size = Pt(45)
p_last.font.name = "Arial"
p_last.font.color.rgb = RGBColor(0, 255, 150)
p_last.alignment = PP_ALIGN.CENTER

out_path = "/home/bozkurt/Desktop/ummuiy/2-darajali/CyberPlace.uz/cdCTF-Pitch-Deck-Pro.pptx"
prs.save(out_path)
print(f"Created amazing pitch deck at {out_path}")
