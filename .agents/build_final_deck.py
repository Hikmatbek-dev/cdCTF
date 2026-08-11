from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import shutil

# Copy screenshots to current dir for easy access
base_artifact_dir = "/home/bozkurt/.gemini/antigravity/brain/b80c17e8-2aa7-44de-8fcc-5bd05055c753/"
try:
    shutil.copyfile(os.path.join(base_artifact_dir, [f for f in os.listdir(base_artifact_dir) if f.startswith('screenshot_hero')][0]), 'screenshot_hero.png')
    shutil.copyfile(os.path.join(base_artifact_dir, [f for f in os.listdir(base_artifact_dir) if f.startswith('screenshot_modules')][0]), 'screenshot_modules.png')
    shutil.copyfile(os.path.join(base_artifact_dir, [f for f in os.listdir(base_artifact_dir) if f.startswith('screenshot_ctf')][0]), 'screenshot_ctf.png')
except Exception as e:
    print(f"Error copying images: {e}")

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

blank_slide_layout = prs.slide_layouts[6] 

BG_COLOR = RGBColor(255, 255, 255)       
TITLE_COLOR = RGBColor(10, 40, 100)     
TEXT_COLOR = RGBColor(40, 50, 60)        
ACCENT_COLOR = RGBColor(10, 100, 220)

def add_slide(title, subtitle, image_path=None, split_layout=True):
    slide = prs.slides.add_slide(blank_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    if image_path and os.path.exists(image_path):
        if split_layout:
            # Side-by-side
            slide.shapes.add_picture(image_path, Inches(8), Inches(2), Inches(7.5), Inches(5))
            text_width = Inches(7)
        else:
            # Big centered image with title on top
            slide.shapes.add_picture(image_path, Inches(3), Inches(3.5), Inches(10), Inches(5))
            text_width = Inches(14)
    else:
        text_width = Inches(14)
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1) if not image_path or split_layout else Inches(0.5), text_width, Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(50)
    p.font.name = "Arial"
    p.font.color.rgb = TITLE_COLOR
    
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3) if not image_path or split_layout else Inches(1.5), text_width, Inches(6))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(26)
        p2.font.name = "Arial"
        p2.font.color.rgb = TEXT_COLOR
    
    line = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(16), Inches(0.2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.color.rgb = ACCENT_COLOR
    
    return slide

add_slide("cdCTF", "O'zbekiston uchun ilg'or kiberxavfsizlik akademiyasi va amaliy CTF platformasi.\nO'rganishdan birinchi ishgacha — bitta platformada.\n\nInvestitsiya maqsadi: $10,000", split_layout=False)

add_slide("01 MUAMMO", "Kiberxavfsizlik kadrlari yetishmaydi:\n\n• O'zbek tilida sifatli, amaliy kiberxavfsizlik ta'limi deyarli yo'q.\n• Xalqaro platformalar qimmat va faqat ingliz tilida.\n• Iqtidorli yoshlar bor, ammo o'rganish bilan ish o'rtasida ko'prik yo'q.")

add_slide("02 YECHIM", "Bitta platforma uch bosqichni bog'laydi — real terminal ko'nikmalaridan birinchi ishgacha.\n\n• O'RGANISH (Akademiya): Noldan ilg'or darajagacha interaktiv darslar.\n• SINOV (CTF): O'rganganlarni amaliyotda tekshirish tizimi.\n• ISHGA CHIQISH (Karyera): HR bazasi va kadrlar yetkazib berish (B2B).")

add_slide("03 MAHSULOT", "Tayyor va ishlayotgan mahsulot. \ncdctf.uz orqali ilk foydalanuvchilar o'qishni boshladi.\n\n• 8 Modul\n• 64 Dars\n• 97 CTF Topshirig'i", "screenshot_hero.png")

add_slide("04 O'QUV DASTURI", "Sistemali o'quv dasturi.\n\nLinux asoslaridan tortib to to'liq tarmoq tizimlarni ekspluatatsiya qilishgacha bo'lgan yo'nalishlar bitta joyda jamlangan.", "screenshot_modules.png")

add_slide("05 AMALIYOT / CTF", "Nazariyani mustahkamlovchi yopiq poligonlar (Sandboxed labs).\n\nO'quvchilar real xakerlik ko'nikmalarini bizning xavfsiz va izolyatsiya qilingan bulutli serverlarimizda sinab ko'rishadi.", "screenshot_ctf.png")

add_slide("06 TEXNOLOGIYA", "Jiddiy muhandislik. Kiberxavfsizlik platformasi o'zi ham xavfsiz bo'lishi shart.\n\n• Stack: React, Vite, Express 5, PostgreSQL, Drizzle ORM\n• Izolyatsiya: Har bir amaliyot alohida Docker konteynerda ishga tushadi\n• Xavfsizlik auditi: Qat'iy tekshiruv va avtomatik testlardan o'tgan")

add_slide("07 BIZNES MODEL", "Biznes daromadini B2B segmentidan shakllantirish rejalashtirilgan:\n\n• Talent Pipeline (Kadrlar): Bitiruvchilarni IT kompaniyalarga yetkazib berish va komissiya olish.\n• Korporativ Ta'lim: Kompaniyalarning o'z dasturchilariga xavfsiz kod yozishni o'rgatuvchi B2B xizmatlar.\n• Premium Sertifikatsiya: Murakkablashtirilgan laboratoriyalar orqali xalqaro standartdagi pullik imtihon tizimi.")

add_slide("08 INVESTITSIYA MAQSADI", "So'ralayotgan sarmoya: $10,000 (10 oylik operatsion xarajatlar uchun)\n\n• Marketing va Jamiyat (40%): O'zbekiston bo'ylab yirik va sovrinli ochiq CTF musobaqalari o'tkazish.\n• Jamoa (40%): Yangi murakkab laboratoriyalar yaratuvchi kiberxavfsizlik muhandislarini jalb qilish.\n• Infratuzilma (20%): Katta hajmdagi bulutli serverlar to'lovlarini qoplash.")

add_slide("09 RAHMAT", "Bir jamoa. Bir missiya.\n\nO'zbekiston uchun ilg'or kiberxavfsizlik ta'limi — barqaror biznes model bilan.\n\nHikmatbek Xudoyberganov\ncdctf.uz | @cdctf", split_layout=False)

prs.save("cdCTF-Pitch-Deck-Final.pptx")
print("Saved cdCTF-Pitch-Deck-Final.pptx")
