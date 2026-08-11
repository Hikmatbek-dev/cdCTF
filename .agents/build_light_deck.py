from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

prs = Presentation()
# Set 16:9 aspect ratio
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

blank_slide_layout = prs.slide_layouts[6] 

# LIGHT THEME COLORS
BG_COLOR = RGBColor(252, 252, 255)       # Almost White
TITLE_COLOR = RGBColor(10, 100, 220)     # Tech Blue
TEXT_COLOR = RGBColor(30, 40, 50)        # Dark Slate for high readability

def add_slide_with_image(title, subtitle, image_path=None):
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Light background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    if image_path and os.path.exists(image_path):
        # Place image on the right side
        slide.shapes.add_picture(image_path, Inches(8.5), Inches(0), Inches(7.5), Inches(9))
        text_width = Inches(7)
    else:
        text_width = Inches(14)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), text_width, Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(54)
    p.font.name = "Arial"
    p.font.color.rgb = TITLE_COLOR
    
    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.0), text_width, Inches(4.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(28)
    p2.font.name = "Arial"
    p2.font.color.rgb = TEXT_COLOR
    
    # Add a minimal colored accent line on the left
    line = slide.shapes.add_shape(
        1, # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(0), Inches(0.2), Inches(9)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TITLE_COLOR
    line.line.color.rgb = TITLE_COLOR
    
    return slide

# 1. Cover
add_slide_with_image("cdCTF", "O'zbekistonning ilg'or kiberxavfsizlik akademiyasi va kadrlar yetkazib berish (Talent Pipeline) tizimi.\n\nInvestitsiya maqsadi: $10,000", "/home/bozkurt/.gemini/antigravity/brain/b80c17e8-2aa7-44de-8fcc-5bd05055c753/light_cover_1786357327114.png")

# 2. Problem
add_slide_with_image("MUAMMO: KADRLAR TAQCHILLIGI", "Bozorda kiberxavfsizlik mutaxassislari keskin yetishmaydi:\n\n• O'zbek tilida sifatli, amaliy ta'lim dasturi yo'q.\n• Yoshlarni real amaliyotga yo'naltiruvchi tizim mavjud emas.\n• HR bo'limlari malakali xavfsizlik muhandislarini topishga qiynalmoqda.", "/home/bozkurt/.gemini/antigravity/brain/b80c17e8-2aa7-44de-8fcc-5bd05055c753/light_problem_1786357337222.png")

# 3. Solution
add_slide_with_image("YECHIM: cdCTF PLATFORMASI", "O'rganishdan tortib ishga joylashishgacha bitta platformada:\n\n• Akademiya: 8 ta modul, 64 ta interaktiv dars (uz/ru/en).\n• CTF Poligon: 97+ izolyatsiya qilingan (Sandboxed) amaliy laboratoriyalar.\n• Karyera Tizimi: Bitiruvchilarni B2B orqali to'g'ridan-to'g'ri ish beruvchilarga bog'lash.", "/home/bozkurt/.gemini/antigravity/brain/b80c17e8-2aa7-44de-8fcc-5bd05055c753/light_solution_1786357348465.png")

# 4. Tech
add_slide_with_image("TEXNOLOGIYA VA XAVFSIZLIK", "Sanoat standartlariga mos va mustahkam arxitektura:\n\n• Zamonaviy Stack: React, Vite, Express 5, PostgreSQL, Drizzle ORM.\n• Xavfsiz Poligonlar: Maxsus Docker konteynerlar yordamida izolyatsiya qilingan server-authoritative laboratoriyalar.\n• Mukammal Himoya: Xavfsizlik auditi va avtomatik testlardan (22 ta to'plam) muvaffaqiyatli o'tgan.")

# 5. Business Model
add_slide_with_image("BIZNES MODEL (B2B)", "Sof va miqyosi kengayuvchi daromad manbalari:\n\n1. Talent Pipeline (Kadrlar):\nKompaniyalarga tasdiqlangan mutaxassislarni tavsiya qilish orqali HR-komissiyasi.\n\n2. Korporativ Ta'lim:\nKompaniyalarning IT jamoalarini xavfsiz kod yozishga o'rgatuvchi yopiq SaaS modullar.\n\n3. Premium Modullar va Sertifikat:\nTalabalar uchun murakkablashtirilgan laboratoriyalar va maxsus imtihonlar sotuvi.")

# 6. Traction
add_slide_with_image("TRACTION VA YUTUQLAR", "Loyiha amalda ishlamoqda va tezkor o'smoqda:\n\n• MVP Holati: cdctf.uz to'liq faol va yuklamalarga chidamli.\n• Kontent: Katta hajmdagi bilimlar bazasi va laboratoriyalar yozib bo'lingan.\n• Multilingual: O'zbek, Rus va Ingliz tillarida muammosiz ishlamoqda.\n• Dastlabki qadam: Foydalanuvchilar qiziqish bildirib, birinchi topshiriqlarni yechishni boshladi.")

# 7. Ask
add_slide_with_image("INVESTITSIYA MAQSADI", "So'ralayotgan sarmoya: $10,000 (10 oylik Operatsion Xarajat)\n\nLoyiha jadal o'sishi uchun sarflanish rejasi:\n• Jamoani Kengaytirish (40%): Yangi va yanada murakkab laboratoriyalar yozish uchun kuchli muhandislar.\n• Marketing va PR (40%): Platformaga o'quvchilarni keng miqyosda jalb qilish.\n• Infratuzilma (20%): Katta yubklamani ko'taruvchi bulutli serverlarni (Cloud) ta'minlash.")

# 8. Contact
slide_last = prs.slides.add_slide(blank_slide_layout)
background = slide_last.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BG_COLOR
txBox_last = slide_last.shapes.add_textbox(Inches(3), Inches(3.5), Inches(10), Inches(4))
tf_last = txBox_last.text_frame
p_last = tf_last.paragraphs[0]
p_last.text = "Bir jamoa. Bir missiya.\n\nHikmatbek Xudoyberganov\ncdctf.uz | @cdctf"
p_last.font.bold = True
p_last.font.size = Pt(45)
p_last.font.name = "Arial"
p_last.font.color.rgb = TITLE_COLOR
p_last.alignment = PP_ALIGN.CENTER

out_path = "/home/bozkurt/Desktop/ummuiy/2-darajali/CyberPlace.uz/cdCTF-Pitch-Deck-LightPro.pptx"
prs.save(out_path)
print(f"Created light mode pitch deck at {out_path}")
